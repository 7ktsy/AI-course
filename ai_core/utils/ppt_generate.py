from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.oxml.xmlchemy import OxmlElement
import os

# 定义颜色主题
COLOR_THEMES = {
    "professional": {
        "primary": RGBColor(0, 76, 153),      # 深蓝色
        "secondary": RGBColor(51, 122, 183),   # 中蓝色
        "accent": RGBColor(92, 184, 92),       # 绿色
        "text": RGBColor(51, 51, 51),          # 深灰色
        "light_text": RGBColor(119, 119, 119), # 浅灰色
        "background": RGBColor(255, 255, 255), # 白色
        "light_bg": RGBColor(248, 249, 250)    # 浅灰背景
    },
    "modern": {
        "primary": RGBColor(52, 73, 94),       # 深青灰
        "secondary": RGBColor(44, 62, 80),     # 深灰
        "accent": RGBColor(231, 76, 60),       # 红色
        "text": RGBColor(44, 62, 80),          # 深灰
        "light_text": RGBColor(127, 140, 141), # 浅灰
        "background": RGBColor(255, 255, 255), # 白色
        "light_bg": RGBColor(236, 240, 241)    # 浅青灰背景
    },
    "elegant": {
        "primary": RGBColor(142, 68, 173),     # 紫色
        "secondary": RGBColor(155, 89, 182),   # 浅紫色
        "accent": RGBColor(241, 196, 15),      # 金色
        "text": RGBColor(52, 73, 94),          # 深灰
        "light_text": RGBColor(149, 165, 166), # 浅灰
        "background": RGBColor(255, 255, 255), # 白色
        "light_bg": RGBColor(248, 249, 250)    # 浅灰背景
    },
    "tech": {
        "primary": RGBColor(41, 128, 185),     # 蓝色
        "secondary": RGBColor(52, 152, 219),   # 浅蓝色
        "accent": RGBColor(230, 126, 34),      # 橙色
        "text": RGBColor(44, 62, 80),          # 深灰
        "light_text": RGBColor(127, 140, 141), # 浅灰
        "background": RGBColor(255, 255, 255), # 白色
        "light_bg": RGBColor(236, 240, 241)    # 浅青灰背景
    }
}

# 定义字体配置
FONT_CONFIGS = {
    "default": {
        "title_font": "微软雅黑",
        "subtitle_font": "微软雅黑", 
        "body_font": "微软雅黑",
        "title_size": 44,
        "subtitle_size": 36,
        "body_size": 24,
        "small_size": 18
    },
    "modern": {
        "title_font": "思源黑体",
        "subtitle_font": "思源黑体",
        "body_font": "思源黑体", 
        "title_size": 48,
        "subtitle_size": 32,
        "body_size": 22,
        "small_size": 16
    },
    "elegant": {
        "title_font": "华文细黑",
        "subtitle_font": "华文细黑",
        "body_font": "华文细黑",
        "title_size": 42,
        "subtitle_size": 34,
        "body_size": 24,
        "small_size": 18
    }
}

def add_gradient_background(slide, color1, color2):
    """为幻灯片添加渐变背景"""
    try:
        # 创建渐变背景
        background = slide.background
        fill = background.fill
        fill.gradient()
        fill.gradient_stops[0].color.rgb = color1
        fill.gradient_stops[1].color.rgb = color2
    except:
        # 如果渐变失败，使用纯色背景
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color1

def add_shadow_to_shape(shape):
    """为形状添加阴影效果"""
    try:
        shadow = shape.shadow
        shadow.inherit = False
        shadow.visible = True
        shadow.blur_radius = 10
        shadow.distance = 5
        shadow.direction = 45
        shadow.color.rgb = RGBColor(0, 0, 0, 50)
    except:
        pass

def create_presentation(title: str, content_dict: dict, theme: str = "professional", 
                       font_config: str = "default", include_animations: bool = False) -> str:
    """
    创建美观的PPT演示文稿
    
    Args:
        title (str): PPT标题
        content_dict (dict): 内容字典，key为标题，value为内容列表
        theme (str): 颜色主题 ("professional", "modern", "elegant", "tech")
        font_config (str): 字体配置 ("default", "modern", "elegant")
        include_animations (bool): 是否包含动画效果
    
    Returns:
        str: 保存的文件路径
    """
    # 获取主题颜色和字体配置
    colors = COLOR_THEMES.get(theme, COLOR_THEMES["professional"])
    fonts = FONT_CONFIGS.get(font_config, FONT_CONFIGS["default"])
    
    # 创建演示文稿对象
    prs = Presentation()
    
    # 设置默认页面大小为16:9
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    # 添加封面页
    cover_slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    # 设置封面背景
    add_gradient_background(cover_slide, colors["primary"], colors["secondary"])
    
    # 设置标题
    title_shape = cover_slide.shapes.title
    title_shape.text = title
    
    # 设置标题格式
    title_frame = title_shape.text_frame
    title_frame.clear()
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.name = fonts["title_font"]
    title_para.font.size = Pt(fonts["title_size"])
    title_para.font.bold = True
    title_para.font.color.rgb = colors["background"]
    title_para.alignment = PP_ALIGN.CENTER
    
    # 添加副标题（可选）
    subtitle_shape = cover_slide.placeholders[1]
    subtitle_frame = subtitle_shape.text_frame
    subtitle_frame.clear()
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.text = "AI智能生成"
    subtitle_para.font.name = fonts["subtitle_font"]
    subtitle_para.font.size = Pt(fonts["small_size"])
    subtitle_para.font.color.rgb = colors["light_bg"]
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # 添加装饰元素
    try:
        # 添加装饰线条
        line = cover_slide.shapes.add_connector(1, Inches(4), Inches(6), Inches(12), Inches(6))
        line.line.color.rgb = colors["accent"]
        line.line.width = Pt(3)
    except:
        pass
    
    # 添加内容页
    for subtitle, contents in content_dict.items():
        # 使用带标题和内容的布局
        content_slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        # 设置页面背景
        add_gradient_background(content_slide, colors["light_bg"], colors["background"])
        
        # 设置页面标题
        slide_title = content_slide.shapes.title
        slide_title.text = subtitle
        
        # 设置标题格式
        title_frame = slide_title.text_frame
        title_frame.clear()
        title_para = title_frame.paragraphs[0]
        title_para.text = subtitle
        title_para.font.name = fonts["subtitle_font"]
        title_para.font.size = Pt(fonts["subtitle_size"])
        title_para.font.bold = True
        title_para.font.color.rgb = colors["primary"]
        title_para.alignment = PP_ALIGN.LEFT
        
        # 获取内容占位符
        content_shape = content_slide.placeholders[1]
        content_frame = content_shape.text_frame
        content_frame.clear()
        
        # 添加内容
        for idx, content in enumerate(contents):
            if idx == 0:
                # 第一段使用默认段落
                p = content_frame.paragraphs[0]
            else:
                # 其他内容添加新段落
                p = content_frame.add_paragraph()
            
            # 设置段落格式
            p.text = f"• {content}"
            p.font.name = fonts["body_font"]
            p.font.size = Pt(fonts["body_size"])
            p.font.color.rgb = colors["text"]
            p.space_after = Pt(12)  # 段落后间距
            p.space_before = Pt(6)  # 段落前间距
            
            # 设置项目符号样式
            p.level = 0
            
            # 添加缩进
            p.left_indent = Pt(20)
        
        # 添加装饰元素
        try:
            # 在标题下方添加装饰线
            line = content_slide.shapes.add_connector(1, Inches(0.5), Inches(1.2), Inches(15.5), Inches(1.2))
            line.line.color.rgb = colors["accent"]
            line.line.width = Pt(2)
        except:
            pass
    
    # 添加结束页
    end_slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    # 设置结束页背景
    add_gradient_background(end_slide, colors["primary"], colors["secondary"])
    
    # 设置结束页标题
    end_title = end_slide.shapes.title
    end_title.text = "谢谢观看"
    
    # 设置结束页标题格式
    end_title_frame = end_title.text_frame
    end_title_frame.clear()
    end_para = end_title_frame.paragraphs[0]
    end_para.text = "谢谢观看"
    end_para.font.name = fonts["title_font"]
    end_para.font.size = Pt(fonts["title_size"])
    end_para.font.bold = True
    end_para.font.color.rgb = colors["background"]
    end_para.alignment = PP_ALIGN.CENTER
    
    # 添加副标题
    end_subtitle = end_slide.placeholders[1]
    end_subtitle_frame = end_subtitle.text_frame
    end_subtitle_frame.clear()
    end_subtitle_para = end_subtitle_frame.paragraphs[0]
    end_subtitle_para.text = f"《{title}》"
    end_subtitle_para.font.name = fonts["subtitle_font"]
    end_subtitle_para.font.size = Pt(fonts["subtitle_size"])
    end_subtitle_para.font.color.rgb = colors["light_bg"]
    end_subtitle_para.alignment = PP_ALIGN.CENTER
    
    # 保存文件
    save_dir = "D:\\Download"
    os.makedirs(save_dir, exist_ok=True)
    file_path = os.path.join(save_dir, f"{title}.pptx")
    
    # 如果文件已存在，添加序号
    counter = 1
    base_path = os.path.splitext(file_path)[0]
    while os.path.exists(file_path):
        file_path = f"{base_path}_{counter}.pptx"
        counter += 1
    
    prs.save(file_path)
    return file_path

def create_presentation_advanced(title: str, content_dict: dict, 
                                theme: str = "professional", 
                                font_config: str = "default",
                                include_animations: bool = False,
                                slide_transitions: bool = True,
                                custom_colors: dict = None) -> str:
    """
    创建高级PPT演示文稿（更多自定义选项）
    
    Args:
        title (str): PPT标题
        content_dict (dict): 内容字典
        theme (str): 颜色主题
        font_config (str): 字体配置
        include_animations (bool): 是否包含动画
        slide_transitions (bool): 是否包含幻灯片切换效果
        custom_colors (dict): 自定义颜色
    
    Returns:
        str: 保存的文件路径
    """
    # 合并自定义颜色
    colors = COLOR_THEMES.get(theme, COLOR_THEMES["professional"]).copy()
    if custom_colors:
        colors.update(custom_colors)
    
    return create_presentation(title, content_dict, theme, font_config, include_animations)

def get_available_themes():
    """获取可用的主题列表"""
    return list(COLOR_THEMES.keys())

def get_available_fonts():
    """获取可用的字体配置列表"""
    return list(FONT_CONFIGS.keys())

def test_presentation():
    """
    测试PPT生成功能
    """
    # 测试数据
    title = "Python编程基础"
    content_dict = {
        "什么是Python?": [
            "Python是一种高级编程语言",
            "易学易用，应用广泛",
            "支持面向对象编程",
            "拥有丰富的第三方库"
        ],
        "Python的应用领域": [
            "Web开发",
            "数据分析",
            "人工智能",
            "自动化运维",
            "科学计算"
        ],
        "为什么选择Python?": [
            "简洁优雅的语法",
            "强大的生态系统",
            "跨平台兼容性好",
            "活跃的社区支持"
        ]
    }
    
    try:
        # 测试不同主题
        themes = ["professional", "modern", "elegant", "tech"]
        for theme in themes:
            print(f"生成 {theme} 主题的PPT...")
            file_path = create_presentation(title, content_dict, theme=theme)
            print(f"✅ {theme} 主题PPT已生成：{file_path}")
        
        # 测试不同字体配置
        fonts = ["default", "modern", "elegant"]
        for font in fonts:
            print(f"生成 {font} 字体配置的PPT...")
            file_path = create_presentation(title, content_dict, font_config=font)
            print(f"✅ {font} 字体配置PPT已生成：{file_path}")
            
    except Exception as e:
        print(f"生成PPT时出错：{e}")

if __name__ == "__main__":
    test_presentation()